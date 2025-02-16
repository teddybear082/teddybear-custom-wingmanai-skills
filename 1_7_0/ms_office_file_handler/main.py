import os
import time
import traceback
from typing import TYPE_CHECKING
from urllib.parse import urlparse
from openpyxl import Workbook, load_workbook
from openpyxl.styles import Font, Alignment
from openpyxl.chart import Reference, BarChart, LineChart, PieChart, ScatterChart, AreaChart
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.util import Inches as pptxInches
from pptx.util import Pt as pptxPt
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.dml.color import RGBColor as pptxRGBColor
from docx import Document
from docx.shared import Inches as docxInches
from docx.shared import Pt as docxPt
from docx.shared import RGBColor as docxRGBColor
from docx.oxml.ns import qn
from api.interface import (
    SettingsConfig,
    SkillConfig,
)
from api.enums import LogType
from services.benchmark import Benchmark
from skills.skill_base import Skill
if TYPE_CHECKING:
    from wingmen.open_ai_wingman import OpenAiWingman

chart_types = {
    "bar": BarChart,
    "line": LineChart,
    "pie": PieChart,
    "scatter": ScatterChart,
    "area": AreaChart
}

class MSOfficeFileHandler(Skill):
    def __init__(
        self,
        config: SkillConfig,
        settings: SettingsConfig,
        wingman: "OpenAiWingman",
    ) -> None:
        super().__init__(config=config, settings=settings, wingman=wingman)

    def get_tools(self) -> list[tuple[str, dict]]:
        tools = [
            (
                "create_excel",
                {
                    "type": "function",
                    "function": {
                        "name": "create_excel",
                        "description": "Creates an Excel file with specified content.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to save the Excel file."
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Name of the sheet.",
                                    "default": "Sheet1"
                                },
                                "data": {
                                    "type": "array",
                                    "items": {
                                        "type": "array",
                                        "items": {
                                            "type": "string"
                                        }
                                    },
                                    "description": "2D array of data to populate the sheet. Cells can also contain formulae, such as '=SUM(1,1)'"
                                },
                                "convert_to_float": {
                                    "type": "boolean",
                                    "description": "Boolean to indicate whether to attempt to format all number like strings in the data as actual numbers in the resulting spreadsheet."
                                },
                                "create_chart": {
                                    "type": "boolean",
                                    "description": "Boolean to indicate whether to create a chart.",
                                },
                                "chart_type": {
                                    "type": "string",
                                    "description": "Type of chart to create, if any. Supported types: 'bar', 'line', 'pie', 'scatter', 'area'.",
                                    "enum": ["bar", "line", "pie", "scatter", "area"]
                                }
                            },
                            "required": ["file_path"]
                        }
                    }
                }
            ),
            (
                "create_ppt",
                {
                    "type": "function",
                    "function": {
                        "name": "create_ppt",
                        "description": "Creates a PowerPoint presentation with specified slides, images, shapes, and title slides.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "slides": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "title": {"type": "string", "description": "The title on the slide"},
                                            "content": {"type": "string", "description": "The rest of the text content of the slide, if any"},
                                            "image_path": {"type": "string", "description": "Path to an image file to be inserted into the slide, if any"},
                                            "shape_type": {"type": "string", "description": "Type of shape to insert into slide, as a type of MSO_SHAPE used in python-pptx, if any"},
                                            "is_title_slide": {"type": "boolean", "description": "Boolean for whether this is a title slide"},
                                            "slide_background_color_rgb": {"type": "array", "items": {"type":"number", "type": "number", "type": "number"}, "description": "Tuple of RGB value for color of background slide color"},
                                            "font": {
                                                "type": "object",
                                                "properties": {
                                                    "name": {"type": "string", "description": "Font name to use for text on slide."},
                                                    "size": {"type": "number", "description": "Font size to use."},
                                                    "color": {"type": "array", "items": {"type":"number", "type": "number", "type": "number"}, "description": "Tuple of RGB value for color of the font."},
                                                },
                                                "description": "Optional font settings for the document."
                                            },
                                        },
                                        "required": ["title"]
                                    }
                                },
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to save the Powerpoint document."
                                }
                            },
                            "required": ["slides", "file_name"]
                        }
                    }
                }
            ),
            (
                "create_docx",
                {
                    "type": "function",
                    "function": {
                        "name": "create_docx",
                        "description": "Creates a Word document with specified content, including optional headers, images, and page breaks.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "file_path": {"type": "string", "description": "Path to save the Word document."},
                                "content": {
                                    "type": "array",
                                    "items": {
                                        "type": "object",
                                        "properties": {
                                            "text": {"type": "string", "description": "The paragraph text."},
                                            "header": {"type": "string", "description": "Optional header for this section."},
                                            "image": {"type": "string", "description": "Optional path to an image to insert."},
                                            "page_break": {"type": "boolean", "description": "Optional page break after this content item.", "default": False},
                                            "list_type": {"type": "string", "description": "Indicates the type of list formatting for the paragraph, if any.", "enum": ["bullet", "numbered"]}
                                        }
                                    },
                                    "description": "List of content items to add to the document."
                                },
                                "font": {
                                    "type": "object",
                                    "properties": {
                                        "name": {"type": "string", "description": "Font name to use."},
                                        "size": {"type": "number", "description": "Font size to use."},
                                        "color": {"type": "array", "items": {"type":"number", "type": "number", "type": "number"}, "description": "Tuple of RGB value for color of the font."},
                                    },
                                    "description": "Optional font settings for the document."
                                }
                            },
                            "required": ["file_path", "content"]
                        }
                    }
                }
            ),
            (
                "read_presentation",
                {
                    "type": "function",
                    "function": {
                        "name": "read_presentation",
                        "description": "Read a PowerPoint presentation and return the text content.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "file_name": {
                                    "type": "string",
                                    "description": "The file name of the presentation to read.",
                                },
                            },
                            "required": ["file_name"],
                        },
                    },
                },
            ),
            (
                "read_document",
                {
                    "type": "function",
                    "function": {
                        "name": "read_document",
                        "description": "Read a Word document and return the text content.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "file_name": {
                                    "type": "string",
                                    "description": "The file name of the document to read.",
                                },
                            },
                            "required": ["file_name"],
                        },
                    },
                },
            ),
            (
                "read_spreadsheet",
                {
                    "type": "function",
                    "function": {
                        "name": "read_spreadsheet",
                        "description": "Read an Excel spreadsheet and return the data.",
                        "parameters": {
                            "type": "object",
                            "properties": {
                                "file_name": {
                                    "type": "string",
                                    "description": "The file name of the spreadsheet to read.",
                                },
                            },
                            "required": ["file_name"],
                        },
                    },
                },
            ),
        ]
        return tools

    async def execute_tool(
        self, tool_name: str, parameters: dict[str, any], benchmark: Benchmark
    ) -> tuple[str, str]:
        function_response = "Error processing command. Please try again."
        instant_response = ""
        if self.settings.debug_mode:
            benchmark.start_snapshot(f"File Manager: {tool_name}")
            await self.printr.print_async(
                f"Executing {tool_name} with parameters: {parameters}",
                color=LogType.INFO,
            )
        # Given many possible issues working with complex files, wrap in try/except block to provide error catching
        try:
            if tool_name == "read_presentation":
                file_name = parameters.get("file_name")
                content = self.read_presentation(file_name)
                function_response = f"Content of {file_name}:\n{content}"
            elif tool_name == "read_document":
                file_name = parameters.get("file_name")
                content = self.read_document(file_name)
                function_response = f"Content of {file_name}:\n{content}"
            elif tool_name == "read_spreadsheet":
                file_name = parameters.get("file_name")
                content = self.read_spreadsheet(file_name)
                function_response = f"Content of {file_name}:\n{content}"
            elif tool_name == "create_excel":
                function_response = await self.create_excel(parameters)
            elif tool_name == "create_ppt":
                function_response = await self.create_ppt(parameters)
            elif tool_name == "create_docx":
                function_response = await self.create_docx(parameters)
        except Exception as e:
            function_response += f"  Error was: {e}; Traceback if any: {traceback.format_exc()}"
        if self.settings.debug_mode:
            await self.printr.print_async(
                f"Output produced by {tool_name}: {function_response}",
                color=LogType.INFO,
            )
            benchmark.finish_snapshot()
        return function_response, instant_response

    # Used for creating charts in Excel
    async def create_chart(self, ws, chart_type: str, data):
        chart_class = chart_types.get(chart_type, BarChart)  # Default to BarChart if type not found
        chart = chart_class()

        # Assuming data starts at the first cell
        data_ref = Reference(ws, min_col=1, min_row=1, max_col=len(data[0]), max_row=len(data))
        chart.add_data(data_ref, titles_from_data=True)

        # Create a new sheet for the chart
        chart_sheet = ws.parent.create_sheet(title=f"{chart_type.capitalize()} Chart")
        chart_sheet.add_chart(chart, "A1")

    # Used for fitting the columns to reasonable widths based on data in Excel
    async def autofit_columns(self, ws, data):
        threshold_width = 30
        for col in ws.columns:
            max_length = 0
            col_letter = get_column_letter(col[0].column)
            for cell in col:
                cell_value = str(cell.value) if cell.value else ""
                max_length = max(max_length, len(cell_value))
            # Apply threshold
            adjusted_width = min(max_length + 2, threshold_width)
            ws.column_dimensions[col_letter].width = adjusted_width

    async def create_excel(self, parameters: dict[str, any]) -> str:
        file_path = parameters["file_path"]
        sheet_name = parameters.get("sheet_name", "Sheet1")
        data = parameters.get("data", [])
        convert_to_float = parameters.get("convert_to_float", True)
        create_chart = parameters.get("create_chart", False)
        chart_type = parameters.get("chart_type", "bar")
        
        wb = Workbook()
        ws = wb.active
        ws.title = sheet_name
        # Apply header formatting
        bold_font = Font(bold=True)
        alignment = Alignment(horizontal='center', vertical='center')
        
        for row_idx, row in enumerate(data, 1):
            new_row = []
            for col_idx, cell in enumerate(row, 1):
                if convert_to_float:
                    try:
                        # Try to convert cell value to float
                        new_cell_value = float(cell)
                    except:
                        # If conversion fails, use the original string value
                        new_cell_value = cell
                else:
                    new_cell_value = cell

                cell_obj = ws.cell(row=row_idx, column=col_idx, value=new_cell_value)

                # Apply formatting to header
                if row_idx == 1:
                    cell_obj.font = bold_font
                    cell_obj.alignment = alignment
\
                # Apply hyperlink formatting if applicable
                if isinstance(new_cell_value, str):
                    parsed_url = urlparse(new_cell_value) 
                    if all([parsed_url.scheme, parsed_url.netloc]):
                        cell_obj.hyperlink = new_cell_value
                        cell_obj.style = "Hyperlink"

                new_row.append(new_cell_value)
            ws.append(new_row)

        # Add filters to all columns
        ws.auto_filter.ref = ws.dimensions

        # Adjust column widths
        await self.autofit_columns(ws, data)

        # Conditionally create a chart
        if create_chart:
            await self.create_chart(ws, chart_type, data)

        # Save final spreadsheet
        wb.save(file_path)
        return f"Excel file created at {file_path}"

    async def create_ppt(self, parameters: dict[str, any]) -> str:
        file_path = parameters["file_path"]
        slides = parameters.get("slides", [])
        prs = Presentation()
        for slide in slides:
            # See https://python-pptx.readthedocs.io/en/latest/user/slides.html for slide layouts (prs.slide_layouts[X])
            if slide.get("is_title_slide"):
                slide_layout = prs.slide_layouts[0]
                ppt_slide = prs.slides.add_slide(slide_layout)
                title = ppt_slide.shapes.title
                subtitle = ppt_slide.placeholders[1]
                title.text = slide["title"]
                subtitle.text = slide.get("content", "")

            else:
                slide_layout = prs.slide_layouts[1]
                ppt_slide = prs.slides.add_slide(slide_layout)
                title = ppt_slide.shapes.title
                title.text = slide["title"]

                # Restrict content to the left side if an image is present
                if slide.get("image_path") and os.path.isfile(slide["image_path"]):
                    left = pptxInches(5)
                    top = pptxInches(1.5)
                    pic = ppt_slide.shapes.add_picture(slide["image_path"], left, top, width=pptxInches(4.5))
                    # Restrict content to the left side
                    content_left = pptxInches(0.5)
                    content_top = pptxInches(1.5)
                    content_width = pptxInches(4.5)
                    content_height = pptxInches(6)
                    content = ppt_slide.shapes.add_textbox(content_left, content_top, content_width, content_height)
                    text_frame = content.text_frame
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.text = slide.get("content", "")
                    
                else:
                    content = ppt_slide.placeholders[1]
                    text_frame = content.text_frame
                    text_frame.word_wrap = True
                    text_frame.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
                    text_frame.text = slide.get("content", "")

                if slide.get("shape_type"):
                    left = pptxInches(1)
                    top = pptxInches(1)
                    width = pptxInches(1)
                    height = pptxInches(1)
                    shape_type = getattr(MSO_SHAPE, slide["shape_type"].upper(), MSO_SHAPE.RECTANGLE)
                    shape = ppt_slide.shapes.add_shape(shape_type, left, top, width, height)

            # Apply slide color if any - https://python-pptx.readthedocs.io/en/latest/dev/analysis/sld-background.html
            possible_color = slide.get("slide_background_color_rgb") 
            if possible_color and len(possible_color) == 3:
                r, g, b = possible_color
                background_fill = ppt_slide.background.fill
                background_fill.solid()
                background_fill.fore_color.rgb = pptxRGBColor(r, g, b)

            # Apply font selection if any
            font_data = slide.get("font")
            if font_data:
                font_name = font_data.get("name")
                font_size = font_data.get("size")
                possible_font_color = font_data.get("color") 

                for shape in ppt_slide.shapes:
                    if not shape.has_text_frame:
                        continue
                    for paragraph in shape.text_frame.paragraphs:
                        for run in paragraph.runs:
                            if font_name:
                                run.font.name = font_name
                            if font_size:
                                run.font.size = pptxPt(font_size)
                            if possible_font_color and len(possible_font_color) == 3:
                                r, g, b = possible_font_color
                                run.font.color.rgb = pptxRGBColor(r, g, b)


        # Save the presentation
        prs.save(file_path)
        return f"Presentation saved at {file_path}"

    async def create_docx(self, parameters: dict[str, any]) -> str:
        file_path = parameters["file_path"]
        content = parameters.get("content", [])
        font_settings = parameters.get("font", {})
        font_name = font_settings.get("name", "Times New Roman")
        font_size = font_settings.get("size", 12)
        possible_font_color = font_settings.get("color", [0,0,0])
        if possible_font_color and len(possible_font_color) == 3:
            r, g, b = possible_font_color
            font_color = docxRGBColor(r,g,b)
        doc = Document()

        # Set default font
        style = doc.styles['Normal']
        font = style.font
        font.name = font_name
        font.size = docxPt(font_size)
        font.color.rgb = font_color

        for item in content:
            # Add header if present
            header = item.get("header")
            if header:
                run = doc.add_heading(header, level=1).runs[0]
                run.font.name = font_name
                run.font.color.rgb = font_color
                run.font.size = docxPt(font_size + 2)
                run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)

            # Add paragraph text and/or image if applicable
            text = item.get("text")
            image = item.get("image")
            list_type = item.get("list_type")
            if text:
                paragraphs = text.split("\n")
                for paragraph_text in paragraphs:
                    if paragraph_text:
                        paragraph = doc.add_paragraph(paragraph_text)
                        if paragraph.runs:
                            run = paragraph.runs[0]
                            run.font.name = font_name
                            run.font.color.rgb = font_color
                            run._element.rPr.rFonts.set(qn('w:eastAsia'), font_name)
                # If text and image, add image in line with text: https://stackoverflow.com/questions/32932230/add-an-image-in-a-specific-position-in-the-document-docx
                if image and run:
                    run.add_picture(image, width=docxInches(5.0))
                # Implement the numbered list and bullet list styles
                # if paragraph and list_type == "bullet":
                    # paragraph.style = 'List Bullet'
                # elif paragraph and list_type == "numbered":
                    # paragraph.style = 'List Number'
                if run and list_type == "bullet":
                    run.style = 'List Bullet'
                elif run and list_type == "numbered":
                    run.style = 'List Number'
            # If just image, and no text, then insert image
            elif image:
                doc.add_picture(image, width=docxInches(5.0))

            # Add page break if specified
            if item.get("page_break", False):
                doc.add_page_break()
        
        doc.save(file_path)
        
        return f"Word document created at {file_path}"

    def read_presentation(self, file_name):
        prs = Presentation(file_name)
        text_runs = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return "\n".join(text_runs)

    def read_spreadsheet(self, file_name):
        wb = load_workbook(file_name)
        ws = wb.active
        data = []
        for row in ws.iter_rows(values_only=True):
            data.append(row)
        return data

    def read_document(self, file_name):
        doc = Document(file_name)
        text = []
        for para in doc.paragraphs:
            text.append(para.text)
        return "\n".join(text)